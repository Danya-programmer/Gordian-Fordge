import csv
import logging
from pathlib import Path
from datetime import datetime
from typing import Optional
from dataclasses import dataclass, asdict

import fitz  # PyMuPDF для PDF
from docx import Document as DocxDocument  # python-docx для DOCX/DOCM
from pptx import Presentation  # python-pptx для PPTX
from openpyxl import load_workbook  # openpyxl для XLSX
import xlrd  # xlrd для XLS

logger = logging.getLogger(__name__)


@dataclass
class FileMetadata:
    """Метаданные файла"""
    author: Optional[str] = None
    created_date: Optional[str] = None
    modified_date: Optional[str] = None
    pages: Optional[int] = None
    slides: Optional[int] = None  # Для PPTX
    sheets: Optional[int] = None  # Для XLSX/XLS
    words: Optional[int] = None
    characters: Optional[int] = None
    title: Optional[str] = None
    subject: Optional[str] = None
    keywords: Optional[str] = None


@dataclass
class ParsedFile:
    """Результат парсинга файла"""
    filename: str
    file_type: str
    file_size: int
    text: str
    metadata: FileMetadata
    
    def to_dict(self) -> dict:
        """Конвертировать в словарь"""
        return {
            "filename": self.filename,
            "file_type": self.file_type,
            "file_size": self.file_size,
            "text": self.text,
            "metadata": asdict(self.metadata)
        }


class FileParser:
    """Парсер файлов различных форматов"""
    
    SUPPORTED_FORMATS = {
        'pdf': ['.pdf'],
        'docx': ['.docx', '.docm'],  # ✅ DOCM добавлен
        'doc': ['.doc'],
        'pptx': ['.pptx'],  # ✅ PPTX добавлен
        'xlsx': ['.xlsx'],
        'xls': ['.xls'],  # ✅ XLS добавлен
        'csv': ['.csv']
    }
    
    def __init__(self):
        self.supported_extensions = set()
        for extensions in self.SUPPORTED_FORMATS.values():
            self.supported_extensions.update(extensions)
    
    def parse(self, file_path: str) -> ParsedFile:
        """
        Парсит файл и возвращает текст + метаданные
        """
        path = Path(file_path)
        
        if not path.exists():
            raise FileNotFoundError(f"Файл не найден: {file_path}")
        
        ext = path.suffix.lower()
        if ext not in self.supported_extensions:
            raise ValueError(f"Формат {ext} не поддерживается. Поддерживаемые: {self.supported_extensions}")
        
        file_type = self._get_file_type(ext)
        
        # Базовая информация о файле
        stat = path.stat()
        file_size = stat.st_size
        modified_date = datetime.fromtimestamp(stat.st_mtime).isoformat()
        
        # Парсим в зависимости от типа
        if file_type == 'pdf':
            text, metadata = self._parse_pdf(file_path)
        elif file_type in ['docx', 'doc']:
            text, metadata = self._parse_docx(file_path)
        elif file_type == 'pptx':
            text, metadata = self._parse_pptx(file_path)
        elif file_type == 'xlsx':
            text, metadata = self._parse_xlsx(file_path)
        elif file_type == 'xls':
            text, metadata = self._parse_xls(file_path)
        elif file_type == 'csv':
            text, metadata = self._parse_csv(file_path)
        else:
            raise ValueError(f"Неизвестный тип файла: {file_type}")
        
        # Добавляем базовые метаданные
        metadata.modified_date = modified_date
        
        # Подсчитываем статистику текста
        if text:
            metadata.words = len(text.split())
            metadata.characters = len(text)
        
        return ParsedFile(
            filename=path.name,
            file_type=file_type,
            file_size=file_size,
            text=text,
            metadata=metadata
        )
    
    def _get_file_type(self, ext: str) -> str:
        """Определяет тип файла по расширению"""
        for file_type, extensions in self.SUPPORTED_FORMATS.items():
            if ext in extensions:
                return file_type
        return 'unknown'
    
    def _parse_pdf(self, file_path: str) -> tuple[str, FileMetadata]:
        """Парсит PDF файл (PyMuPDF — быстро)"""
        logger.info(f"📄 Парсинг PDF: {file_path}")
        
        doc = fitz.open(file_path)
        metadata = FileMetadata()
        
        # Метаданные PDF
        pdf_metadata = doc.metadata
        metadata.author = pdf_metadata.get('author')
        metadata.created_date = pdf_metadata.get('creationDate')
        metadata.title = pdf_metadata.get('title')
        metadata.subject = pdf_metadata.get('subject')
        metadata.keywords = pdf_metadata.get('keywords')
        metadata.pages = len(doc)
        
        # Извлекаем текст с сортировкой (улучшает качество для русских PDF)
        text_parts = []
        for page in doc:
            text = page.get_text("text", sort=True)
            if text.strip():
                text_parts.append(text)
        
        text = '\n\n'.join(text_parts)
        doc.close()
        
        return text, metadata
    
    def _parse_docx(self, file_path: str) -> tuple[str, FileMetadata]:
        """Парсит DOCX/DOCM файл"""
        logger.info(f"📝 Парсинг DOCX/DOCM: {file_path}")
        
        doc = DocxDocument(file_path)
        metadata = FileMetadata()
        
        # Метаданные
        core_props = doc.core_properties
        metadata.author = core_props.author
        metadata.created_date = core_props.created.isoformat() if core_props.created else None
        metadata.modified_date = core_props.modified.isoformat() if core_props.modified else None
        metadata.title = core_props.title
        metadata.subject = core_props.subject
        metadata.keywords = core_props.keywords
        
        # Извлекаем текст из параграфов
        text_parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        # ✅ Также извлекаем текст из таблиц
        for table in doc.tables:
            for row in table.rows:
                row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    text_parts.append(row_text)
        
        text = '\n\n'.join(text_parts)
        
        return text, metadata
    
    def _parse_pptx(self, file_path: str) -> tuple[str, FileMetadata]:
        """Парсит PPTX файл (PowerPoint)"""
        logger.info(f"📊 Парсинг PPTX: {file_path}")
        
        prs = Presentation(file_path)
        metadata = FileMetadata()
        
        # Метаданные
        metadata.author = prs.core_properties.author
        metadata.created_date = prs.core_properties.created.isoformat() if prs.core_properties.created else None
        metadata.modified_date = prs.core_properties.modified.isoformat() if prs.core_properties.modified else None
        metadata.title = prs.core_properties.title
        metadata.subject = prs.core_properties.subject
        metadata.keywords = prs.core_properties.keywords
        metadata.slides = len(prs.slides)
        
        # Извлекаем текст из всех слайдов
        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = [f"=== Слайд {slide_num} ==="]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text)
                
                # ✅ Также извлекаем текст из таблиц
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                        if row_text:
                            slide_texts.append(row_text)
                
                # ✅ Извлекаем текст из групп
                if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                    for child_shape in shape.shapes:
                        if hasattr(child_shape, "text") and child_shape.text.strip():
                            slide_texts.append(child_shape.text)
            
            if len(slide_texts) > 1:  # Есть текст кроме заголовка слайда
                text_parts.append('\n'.join(slide_texts))
        
        text = '\n\n'.join(text_parts)
        
        return text, metadata
    
    def _parse_xlsx(self, file_path: str) -> tuple[str, FileMetadata]:
        """Парсит XLSX файл"""
        logger.info(f"📊 Парсинг XLSX: {file_path}")
        
        wb = load_workbook(file_path, read_only=True, data_only=True)
        metadata = FileMetadata()
        
        # Метаданные
        metadata.author = wb.properties.creator
        metadata.created_date = wb.properties.created
        metadata.modified_date = wb.properties.modified
        metadata.title = wb.properties.title
        metadata.subject = wb.properties.subject
        metadata.keywords = wb.properties.keywords
        metadata.sheets = len(wb.sheetnames)
        
        # Извлекаем текст из всех листов
        text_parts = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text_parts.append(f"=== Лист: {sheet_name} ===")
            
            for row in sheet.iter_rows(values_only=True):
                row_text = ' | '.join(str(cell) if cell is not None else '' for cell in row)
                if row_text.strip('| '):
                    text_parts.append(row_text)
        
        wb.close()
        text = '\n'.join(text_parts)
        
        return text, metadata
    
    def _parse_xls(self, file_path: str) -> tuple[str, FileMetadata]:
        """Парсит XLS файл (старый Excel 97-2003)"""
        logger.info(f"📊 Парсинг XLS: {file_path}")
        
        wb = xlrd.open_workbook(file_path)
        metadata = FileMetadata()
        
        # ✅ Безопасно читаем только то, что точно есть в xlrd
        try:
            metadata.author = getattr(wb, 'user_name', None)
        except Exception:
            pass
        
        metadata.sheets = wb.nsheets
        
        # ❌ УДАЛЕНО (этих атрибутов НЕТ в xlrd):
        # wb.creator
        # wb.created_date
        # wb.modified_date
        # wb.title
        # wb.subject
        # wb.keywords
        
        # Извлекаем текст
        text_parts = []
        for sheet_idx in range(wb.nsheets):
            sheet = wb.sheet_by_index(sheet_idx)
            text_parts.append(f"=== Лист: {sheet.name} ===")
            
            for row_idx in range(sheet.nrows):
                row = sheet.row_values(row_idx)
                row_text = ' | '.join(
                    str(cell) if cell not in (None, '', 0) else '' 
                    for cell in row
                )
                if row_text.strip('| '):
                    text_parts.append(row_text)
        
        text = '\n'.join(text_parts)
        
        return text, metadata
    
    def _parse_csv(self, file_path: str) -> tuple[str, FileMetadata]:
        """Парсит CSV файл"""
        logger.info(f"📋 Парсинг CSV: {file_path}")
        
        metadata = FileMetadata()
        text_parts = []
        
        # Пробуем разные кодировки
        encodings = ['utf-8', 'cp1251', 'latin-1']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    reader = csv.reader(f)
                    for row in reader:
                        text_parts.append(' | '.join(row))
                break
            except UnicodeDecodeError:
                continue
        
        text = '\n'.join(text_parts)
        
        return text, metadata


# Глобальный экземпляр парсера
file_parser = FileParser()


def parse_file(file_path: str) -> dict:
    """
    Удобная функция для парсинга файла
    """
    parsed = file_parser.parse(file_path)
    return parsed.to_dict()